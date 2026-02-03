import streamlit as st
import cv2
import numpy as np
import easyocr
from PIL import Image
import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu

# SSL 인증서 오류 방지
import ssl
ssl._create_default_https_context = ssl._create_unverified_context

def process_image(image_bytes):
    nparr = np.frombuffer(image_bytes, np.uint8)
    img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
    h, w, _ = img.shape # 이미지의 높이와 너비 추출
    
    reader = easyocr.Reader(['ko', 'en'])
    results = reader.readtext(img)
    
    mask = np.zeros(img.shape[:2], dtype="uint8")
    for (bbox, text, prob) in results:
        top_left = tuple(map(int, bbox[0]))
        bottom_right = tuple(map(int, bbox[2]))
        cv2.rectangle(mask, top_left, bottom_right, 255, -1)
    
    clean_img = cv2.inpaint(img, mask, inpaintRadius=7, flags=cv2.INPAINT_TELEA)
    clean_img_rgb = cv2.cvtColor(clean_img, cv2.COLOR_BGR2RGB)
    
    return clean_img_rgb, results, w, h

def create_pptx(clean_image_array, ocr_results, width_px, height_px):
    prs = Presentation()
    
    # 1. 슬라이드 크기를 이미지 크기와 동일하게 설정 (96 DPI 기준 변환)
    # 1 pixel = 9525 EMU (English Metric Units)
    prs.slide_width = Emu(width_px * 9525)
    prs.slide_height = Emu(height_px * 9525)
    
    # --- 슬라이드 1: 텍스트 제거 이미지 ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    img_pil = Image.fromarray(clean_image_array)
    img_io = io.BytesIO()
    img_pil.save(img_io, format='PNG')
    img_io.seek(0)
    slide1.shapes.add_picture(img_io, 0, 0, width=prs.slide_width, height=prs.slide_height)
    
    # --- 슬라이드 2: 추출된 텍스트 (위치 및 크기 반영) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    for (bbox, text, prob) in ocr_results:
        # 좌표 정보 (픽셀 단위)
        x = bbox[0][0]
        y = bbox[0][1]
        w_box = bbox[2][0] - bbox[0][0]
        h_box = bbox[2][1] - bbox[0][1]
        
        # 텍스트 박스 추가 (좌표 변환)
        txBox = slide2.shapes.add_textbox(Emu(x * 9525), Emu(y * 9525), Emu(w_box * 9525), Emu(h_box * 9525))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = text
        
        # 폰트 크기 계산 (상자 높이의 약 70~80% 정도로 자동 조절)
        # 포인트(Pt) 단위로 변환 필요 (1 pixel ≈ 0.75 point)
        font_size = h_box * 0.75
        p.font.size = Pt(max(6, font_size)) # 최소 6pt 보장
        
        # 굵기 처리 (확률이 높거나 특정 조건일 때 굵게 가능하나, 여기선 기본 반영)
        p.font.bold = True if prob > 0.5 else False

    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io

# --- Streamlit UI ---
st.set_page_config(page_title="고급 이미지-PPT 변환기", layout="centered")
st.title("🎯 맞춤형 PPT 생성기")
st.write("이미지의 크기, 텍스트 위치, 크기를 그대로 보존하여 PPT를 제작합니다.")

uploaded_file = st.file_uploader("이미지를 업로드하세요", type=["jpg", "png", "jpeg"])

if uploaded_file:
    file_bytes = uploaded_file.read()
    
    if st.button("🚀 분석 및 PPT 생성"):
        with st.spinner('이미지 구조를 분석하고 있습니다...'):
            clean_img, ocr_results, w, h = process_image(file_bytes)
            ppt_file = create_pptx(clean_img, ocr_results, w, h)
            
            st.success(f"분석 완료! (이미지 사이즈: {w}x{h} px)")
            
            # 결과물 다운로드
            st.download_button(
                label="📊 맞춤형 PPTX 다운로드",
                data=ppt_file,
                file_name="custom_layout.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
            
            st.image(clean_img, caption="텍스트가 제거된 이미지 (1페이지용)", use_container_width=True)

